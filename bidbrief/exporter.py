"""Export results: Excel / Word / PPT / Markdown reports + per-file raw JSON."""
import json
import os
from datetime import datetime

from .extractor import CATEGORIES

# Supported export formats and their dispatch.
FORMATS = ("xlsx", "docx", "pptx", "md")


def _item_text(cat, item):
    if cat == "时间节点":
        return f"{item.get('事项', '')}：{item.get('时间', '')}".strip("：")
    if cat == "评分标准":
        parts = [item.get("评分项", "")]
        if item.get("分值"):
            parts.append(f"分值：{item['分值']}")
        if item.get("说明"):
            parts.append(item["说明"])
        return "；".join(p for p in parts if p)
    return item.get("内容", "")


def _page_label(item):
    p = item.get("页码")
    if p is None:
        return "—"
    if isinstance(p, str):
        return p
    return str(p)


def export_all(file_results, out_dir, formats=None):
    """Export all results into out_dir in the requested formats.

    formats is a subset of FORMATS (defaults to all); raw JSON is always
    written. Returns {format: path} for the generated files.
    """
    if formats is None:
        formats = FORMATS
    formats = [f for f in formats if f in FORMATS]

    os.makedirs(out_dir, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    raw_dir = os.path.join(out_dir, "raw")
    os.makedirs(raw_dir, exist_ok=True)

    for fr in file_results:
        base = os.path.splitext(os.path.basename(fr["path"]))[0]
        with open(os.path.join(raw_dir, f"{base}.json"), "w", encoding="utf-8") as f:
            json.dump(fr, f, ensure_ascii=False, indent=2)

    paths = {"raw_dir": raw_dir}
    if "xlsx" in formats:
        p = os.path.join(out_dir, f"要点报告_{ts}.xlsx")
        _export_excel(file_results, p)
        paths["excel"] = p
    if "docx" in formats:
        p = os.path.join(out_dir, f"要点报告_{ts}.docx")
        _export_docx(file_results, p)
        paths["word"] = p
    if "pptx" in formats:
        p = os.path.join(out_dir, f"要点汇报_{ts}.pptx")
        _export_pptx(file_results, p)
        paths["ppt"] = p
    if "md" in formats:
        p = os.path.join(out_dir, f"核对清单_{ts}.md")
        _export_markdown(file_results, p)
        paths["markdown"] = p
    return paths


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------
def _export_excel(file_results, path):
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="4472C4")
    red_fill = PatternFill("solid", fgColor="FFC7CE")
    wrap = Alignment(wrap_text=True, vertical="top")

    # ---- Sheet 1: all extracted items ----
    ws = wb.active
    ws.title = "要点总表"
    headers = ["序号", "类别", "内容", "原文摘录", "页码", "来源文件"]
    ws.append(headers)
    for c, _ in enumerate(headers, 1):
        cell = ws.cell(1, c)
        cell.font = head_font
        cell.fill = head_fill
    row = 0
    for fr in file_results:
        fname = os.path.basename(fr["path"])
        for cat in CATEGORIES:
            for item in fr.get("categories", {}).get(cat, []):
                row += 1
                ws.append([
                    row, cat, _item_text(cat, item), item.get("原文摘录", ""),
                    _page_label(item), fname,
                ])
                r = ws.max_row
                for c in range(1, 7):
                    ws.cell(r, c).alignment = wrap
                if cat == "废标项":
                    for c in range(1, 7):
                        ws.cell(r, c).fill = red_fill
    widths = [6, 10, 60, 50, 10, 30]
    for c, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(c)].width = w
    ws.freeze_panes = "A2"

    # ---- Sheet 2: disqualification checklist ----
    ws2 = wb.create_sheet("废标核对单")
    headers2 = ["序号", "废标条款", "原文摘录", "页码", "来源文件", "是否满足(填:是/否)"]
    ws2.append(headers2)
    for c, _ in enumerate(headers2, 1):
        cell = ws2.cell(1, c)
        cell.font = head_font
        cell.fill = PatternFill("solid", fgColor="C00000")
    n = 0
    for fr in file_results:
        fname = os.path.basename(fr["path"])
        for item in fr.get("categories", {}).get("废标项", []):
            n += 1
            ws2.append([
                n, item.get("内容", ""), item.get("原文摘录", ""),
                _page_label(item), fname, "",
            ])
            for c in range(1, 7):
                ws2.cell(ws2.max_row, c).alignment = wrap
    for c, w in enumerate([6, 60, 50, 10, 30, 18], 1):
        ws2.column_dimensions[get_column_letter(c)].width = w
    ws2.freeze_panes = "A2"

    # ---- Sheet 3: per-file overview ----
    ws3 = wb.create_sheet("文件概览")
    headers3 = ["文件名", "状态", "总页数", "要点条数"]
    ws3.append(headers3)
    for c, _ in enumerate(headers3, 1):
        cell = ws3.cell(1, c)
        cell.font = head_font
        cell.fill = head_fill
    for fr in file_results:
        cnt = sum(len(v) for v in fr.get("categories", {}).values())
        pages = fr.get("page_count")
        ws3.append([
            os.path.basename(fr["path"]), fr.get("status", ""),
            pages if pages is not None else "—", cnt,
        ])
    for c, w in enumerate([40, 22, 10, 10], 1):
        ws3.column_dimensions[get_column_letter(c)].width = w

    wb.save(path)


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------
def _docx_apply_ea_font(doc):
    """Point the common styles at a Chinese-capable East Asian font."""
    from docx.oxml.ns import qn

    for name in ("Normal", "Title", "Heading 1", "Heading 2"):
        try:
            style = doc.styles[name]
            style.font.name = style.font.name or "Calibri"  # ensure rPr/rFonts exist
            style.element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        except Exception:
            pass


def _docx_table(doc, headers, rows):
    from docx.shared import Cm

    t = doc.add_table(rows=1, cols=len(headers))
    t.style = "Table Grid"
    for i, h in enumerate(headers):
        run = t.rows[0].cells[i].paragraphs[0].add_run(h)
        run.bold = True
    for row in rows:
        cells = t.add_row().cells
        for i, v in enumerate(row):
            cells[i].text = str(v)
    return t


def _export_docx(file_results, path):
    from docx import Document
    from docx.shared import Cm, Pt

    doc = Document()
    _docx_apply_ea_font(doc)
    try:
        doc.styles["Normal"].font.size = Pt(10.5)
    except Exception:
        pass

    doc.add_heading("招标文件要点核对清单", level=0)
    doc.add_paragraph(
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}　　文件数：{len(file_results)}"
    )
    doc.add_paragraph("提示：本清单由 AI 辅助抽取，页码可跳转核对，投标决策前请逐条核对原文。")

    doc.add_heading("一、文件概览", level=1)
    rows = []
    for fr in file_results:
        cnt = sum(len(v) for v in fr.get("categories", {}).values())
        pages = fr.get("page_count")
        rows.append([
            os.path.basename(fr["path"]), fr.get("status", ""),
            pages if pages is not None else "—", cnt,
        ])
    _docx_table(doc, ["文件名", "状态", "页数", "要点数"], rows)

    doc.add_heading("二、废标项核对单（一票否决，逐条确认）", level=1)
    rows = []
    n = 0
    for fr in file_results:
        fname = os.path.basename(fr["path"])
        for item in fr.get("categories", {}).get("废标项", []):
            n += 1
            rows.append([
                n, _item_text("废标项", item), item.get("原文摘录", ""),
                _page_label(item), fname, "☐",
            ])
    if rows:
        _docx_table(doc, ["序号", "废标条款", "原文摘录", "页码", "来源文件", "核对"], rows)
    else:
        doc.add_paragraph("（未抽取到废标项——请注意核对是否遗漏）")

    doc.add_heading("三、分类要点", level=1)
    for fr in file_results:
        cats = fr.get("categories", {})
        if not any(cats.get(c) for c in CATEGORIES):
            continue
        doc.add_heading(os.path.basename(fr["path"]), level=2)
        for cat in CATEGORIES:
            items = cats.get(cat, [])
            if not items:
                continue
            doc.add_heading(cat, level=3)
            rows = [
                [i, _item_text(cat, item), item.get("原文摘录", ""), _page_label(item)]
                for i, item in enumerate(items, 1)
            ]
            _docx_table(doc, ["序号", "内容", "原文摘录", "页码"], rows)

    doc.save(path)


# ---------------------------------------------------------------------------
# PowerPoint (condensed summary for review meetings; full detail lives in
# the Excel/Word reports)
# ---------------------------------------------------------------------------
_PPT_MAX_BULLETS = 7  # bullets per slide, kept low for readability


def _ppt_add_slide(prs, title):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.name = "微软雅黑"
    return slide


def _ppt_bullets(slide, lines, size=15):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.name = "微软雅黑"
        p.space_after = Pt(8)


def _ppt_paginate(prs, title, lines, max_lines=_PPT_MAX_BULLETS):
    """Lay lines out over as many slides as needed, with (x/n) titles."""
    total = max(1, (len(lines) + max_lines - 1) // max_lines)
    for page in range(total):
        part = lines[page * max_lines : (page + 1) * max_lines]
        suffix = f"（{page + 1}/{total}）" if total > 1 else ""
        slide = _ppt_add_slide(prs, title + suffix)
        _ppt_bullets(slide, part)


def _export_pptx(file_results, path):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "招标文件要点提取汇报"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}　｜　文件数：{len(file_results)}"
    )
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(16)
            run.font.name = "微软雅黑"

    # Overview slide
    lines = []
    for fr in file_results:
        cnt = sum(len(v) for v in fr.get("categories", {}).values())
        pages = fr.get("page_count")
        pages_txt = str(pages) if pages is not None else "—"
        lines.append(f"{os.path.basename(fr['path'])}：{fr.get('status', '')}，"
                     f"{pages_txt} 页，{cnt} 条要点")
    _ppt_paginate(prs, "文件概览", lines)

    # Disqualification checklist slides
    lines = []
    for fr in file_results:
        for item in fr.get("categories", {}).get("废标项", []):
            lines.append(f"☐ {_item_text('废标项', item)}（第{_page_label(item)}页）")
    if lines:
        _ppt_paginate(prs, "废标项核对单（一票否决）", lines)
    else:
        slide = _ppt_add_slide(prs, "废标项核对单（一票否决）")
        _ppt_bullets(slide, ["（未抽取到废标项——请注意核对是否遗漏）"])

    # Per-file category slides
    for fr in file_results:
        cats = fr.get("categories", {})
        if not any(cats.get(c) for c in CATEGORIES):
            continue
        fname = os.path.splitext(os.path.basename(fr["path"]))[0]
        for cat in CATEGORIES:
            items = cats.get(cat, [])
            if not items:
                continue
            lines = [f"{_item_text(cat, item)}（第{_page_label(item)}页）" for item in items]
            _ppt_paginate(prs, f"{fname} — {cat}", lines)

    prs.save(path)


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------
def _export_markdown(file_results, path):
    lines = []
    lines.append("# 招标文件要点核对清单")
    lines.append("")
    lines.append(f"- 生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    lines.append(f"- 文件数：{len(file_results)}")
    lines.append("")
    lines.append("> 提示：本清单由 AI 辅助抽取，页码可供跳转核对，投标决策前请逐条核对原文。")
    lines.append("")

    lines.append("## 一、文件概览")
    lines.append("")
    lines.append("| 文件 | 状态 | 页数 | 要点数 |")
    lines.append("| --- | --- | --- | --- |")
    for fr in file_results:
        cnt = sum(len(v) for v in fr.get("categories", {}).values())
        pages = fr.get("page_count")
        lines.append(
            f"| {os.path.basename(fr['path'])} | {fr.get('status', '')} | "
            f"{pages if pages is not None else '—'} | {cnt} |"
        )
    lines.append("")

    lines.append("## 二、废标项核对单（一票否决，逐条确认）")
    lines.append("")
    n = 0
    for fr in file_results:
        for item in fr.get("categories", {}).get("废标项", []):
            n += 1
            lines.append(f"- [ ] **{n}.** {_item_text('废标项', item)}（第{_page_label(item)}页）")
            if item.get("原文摘录"):
                lines.append(f"  > {item['原文摘录']}")
    if n == 0:
        lines.append("（未抽取到废标项——请注意核对是否遗漏）")
    lines.append("")

    lines.append("## 三、分类要点")
    lines.append("")
    for fr in file_results:
        cats = fr.get("categories", {})
        if not any(cats.get(c) for c in CATEGORIES):
            continue
        lines.append(f"### {os.path.basename(fr['path'])}")
        lines.append("")
        for cat in CATEGORIES:
            items = cats.get(cat, [])
            if not items:
                continue
            lines.append(f"#### {cat}")
            lines.append("")
            lines.append("| 序号 | 内容 | 原文摘录 | 页码 |")
            lines.append("| --- | --- | --- | --- |")
            for i, item in enumerate(items, 1):
                quote = item.get("原文摘录", "").replace("|", "\\|")
                content = _item_text(cat, item).replace("|", "\\|").replace("\n", " ")
                lines.append(f"| {i} | {content} | {quote} | {_page_label(item)} |")
            lines.append("")

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
