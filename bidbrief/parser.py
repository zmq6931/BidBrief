"""Parse input files into text.

Supported: PDF, DOCX, XLSX, XLS, PPTX, PPT, TXT, MD.
PDFs keep per-page text for page-level citation; the other formats carry no
page information (page column shows a dash).
"""
import os

# Extensions handled by parse_file (keep in sync with scanner.py).
INPUT_EXTENSIONS = (".pdf", ".docx", ".xlsx", ".xls", ".ppt", ".pptx", ".txt", ".md")


def parse_file(path):
    """Parse a single file.

    Returns a dict:
      ok         whether text was successfully extracted
      pages      list of per-page texts for PDFs; None for all other formats
      text       full document text
      is_scanned PDF looks like a pure scan (no text layer)
      error      failure reason (set when ok is False)
    """
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return _parse_pdf(path)
        if ext == ".docx":
            return _parse_docx(path)
        if ext == ".xlsx":
            return _parse_xlsx(path)
        if ext == ".xls":
            return _parse_xls(path)
        if ext == ".pptx":
            return _parse_pptx(path)
        if ext == ".ppt":
            return _parse_ppt(path)
        if ext in (".txt", ".md"):
            return _parse_text(path)
        return _fail(f"不支持的文件格式: {ext}")
    except Exception as e:
        return _fail(f"解析失败: {e}")


def _fail(msg, **extra):
    d = {"ok": False, "pages": None, "text": "", "is_scanned": False, "error": msg}
    d.update(extra)
    return d


def _ok(text):
    return {"ok": True, "pages": None, "text": text, "is_scanned": False, "error": None}


def _cell_str(v):
    """Normalize a spreadsheet cell value to a compact string."""
    if v is None:
        return ""
    if isinstance(v, float) and v == int(v):
        v = int(v)  # 5.0 -> 5
    return str(v).strip()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def _parse_pdf(path):
    import pymupdf  # PyMuPDF

    doc = pymupdf.open(path)
    if doc.needs_pass:
        if not doc.authenticate(""):
            doc.close()
            return _fail("PDF 已加密，无法读取（部分投标平台下载的文件带 CA 加密）")
    pages = []
    for page in doc:
        try:
            pages.append(page.get_text("text") or "")
        except Exception:
            pages.append("")
    doc.close()

    total = sum(len(p.strip()) for p in pages)
    n = max(len(pages), 1)
    if total / n < 50:
        # Almost no text layer: treat the whole file as a scan.
        return {
            "ok": True,
            "pages": pages,
            "text": "\n".join(pages),
            "is_scanned": True,
            "error": None,
        }
    return {"ok": True, "pages": pages, "text": "\n".join(pages), "is_scanned": False, "error": None}


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------
def _parse_docx(path):
    from docx import Document

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs]
    for t in doc.tables:  # scoring criteria often live inside tables
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(x for x in cells if x)
            if line:
                parts.append(line)
    text = "\n".join(x for x in parts if x.strip())
    if not text.strip():
        return _fail("DOCX 中未提取到文本（可能为图片内容）")
    return _ok(text)


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------
def _parse_xlsx(path):
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"【工作表：{ws.title}】")
        for row in ws.iter_rows(values_only=True):
            line = " | ".join(x for x in (_cell_str(v) for v in row) if x)
            if line:
                parts.append(line)
    wb.close()
    text = "\n".join(x for x in parts if x.strip())
    if not text.strip():
        return _fail("XLSX 中未提取到文本（可能为纯图片/图表）")
    return _ok(text)


def _parse_xls(path):
    import xlrd

    book = xlrd.open_workbook(path)
    parts = []
    for sheet in book.sheets():
        parts.append(f"【工作表：{sheet.name}】")
        for row in sheet.get_rows():
            vals = []
            for cell in row:
                if cell.ctype in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
                    continue
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        dt = xlrd.xldate_as_datetime(cell.value, book.datemode)
                        vals.append(dt.strftime("%Y-%m-%d"))
                    except Exception:
                        vals.append(_cell_str(cell.value))
                else:
                    vals.append(_cell_str(cell.value))
            line = " | ".join(x for x in vals if x)
            if line:
                parts.append(line)
    text = "\n".join(x for x in parts if x.strip())
    if not text.strip():
        return _fail("XLS 中未提取到文本（可能为纯图片/图表）")
    return _ok(text)


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------
def _parse_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    line = " | ".join(x for x in (c.text.strip() for c in row.cells) if x)
                    if line:
                        texts.append(line)
        if texts:
            parts.append(f"【幻灯片 {i}】\n" + "\n".join(texts))
    text = "\n\n".join(parts)
    if not text.strip():
        return _fail("PPTX 中未提取到文本（可能为纯图片）")
    return _ok(text)


def _parse_ppt(path):
    # Legacy binary .ppt: no pure-Python reader; go through the local
    # PowerPoint COM interface (Windows + Office required).
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return _fail("读取 .ppt 需要 pywin32（pip install pywin32），且本机需安装 PowerPoint")
    pythoncom.CoInitialize()
    app = None
    pres = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        # Open(FileName, ReadOnly, Untitled, WithWindow) - keep it headless.
        pres = app.Presentations.Open(os.path.abspath(path), True, False, False)
        parts = []
        for i in range(1, pres.Slides.Count + 1):
            texts = []
            for shape in pres.Slides(i).Shapes:
                try:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        texts.append(shape.TextFrame.TextRange.Text)
                except Exception:
                    pass
            if texts:
                parts.append(f"【幻灯片 {i}】\n" + "\n".join(texts))
        text = "\n\n".join(parts)
        if not text.strip():
            return _fail("PPT 中未提取到文本（可能为纯图片）")
        return _ok(text)
    except Exception as e:
        return _fail(f"读取 .ppt 失败（本机需安装 PowerPoint；或先另存为 .pptx）：{e}")
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        pythoncom.CoUninitialize()


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------
def _parse_text(path):
    text = ""
    for enc in ("utf-8", "gb18030"):
        try:
            with open(path, "r", encoding=enc) as f:
                text = f.read()
            break
        except UnicodeDecodeError:
            continue
    if not text.strip():
        return _fail("文件为空或无法解码")
    return _ok(text)
